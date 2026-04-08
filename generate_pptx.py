"""Generate GenBI Solution Presentation PowerPoint"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BG = RGBColor(0x23, 0x2F, 0x3E)       # AWS dark
BLUE = RGBColor(0x00, 0x73, 0xBB)           # AWS blue
LIGHT_BLUE = RGBColor(0x60, 0xA5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MEDIUM_GRAY = RGBColor(0x88, 0x88, 0x88)
ORANGE = RGBColor(0xFF, 0x99, 0x00)          # AWS orange
RED_ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)
YELLOW_ACCENT = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(text_frame, text, font_size=16, color=WHITE, bold=False, space_before=6, bullet=False):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Segoe UI'
    p.space_before = Pt(space_before)
    if bullet:
        p.level = 0
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color, text='', font_size=14, font_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = 'Segoe UI'
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def add_section_header(slide, text, top=0.3):
    add_text(slide, 0.8, top, 11, 0.6, text, font_size=28, color=WHITE, bold=True)
    # Orange underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top + 0.65), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 1.5, 11, 1.2, 'GenBI Analytics Platform', font_size=48, color=WHITE, bold=True)
add_text(slide, 0.8, 2.7, 11, 0.8, 'AI-Powered Business Intelligence with Data Lineage', font_size=28, color=LIGHT_BLUE)
add_text(slide, 0.8, 3.8, 11, 0.5, 'Natural Language Q&A  |  5 Interactive Dashboards  |  End-to-End Pipeline Tracing', font_size=18, color=MEDIUM_GRAY)

# Service badges
services = ['Amazon Bedrock', 'Amazon Redshift', 'Amazon QuickSight', 'AWS Glue', 'Amazon S3', 'AWS Lambda']
for i, svc in enumerate(services):
    add_rounded_rect(slide, 0.8 + i * 2.05, 5.0, 1.9, 0.45, RGBColor(0x1A, 0x3A, 0x5F), svc, font_size=11, font_color=LIGHT_BLUE)

add_text(slide, 0.8, 6.2, 5, 0.4, 'ABC Restaurant Group  |  200 Stores  |  Hong Kong', font_size=14, color=MEDIUM_GRAY)
add_text(slide, 9.5, 6.2, 3.5, 0.4, 'Powered by Amazon Web Services', font_size=14, color=ORANGE, alignment=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'The Problem')

problems = [
    ('Too Many Dashboards', 'Organizations maintain dozens of dashboards across sales, operations, finance, and customer analytics. Finding the right dashboard for a specific question wastes time and creates frustration.', RED_ACCENT),
    ('Numbers Without Context', 'A dashboard shows "Total Orders: 7,462,065" but doesn\'t explain what\'s included, what\'s excluded, or how it was calculated. Is it line items or distinct transactions?', YELLOW_ACCENT),
    ('Opaque Data Pipelines', 'When a number looks wrong, tracing it from dashboard to data warehouse to ETL to raw source requires tribal knowledge that lives in people\'s heads, not in the tool.', ORANGE),
]

for i, (title, desc, color) in enumerate(problems):
    y = 1.3 + i * 2.0
    # Color bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(0.08), Inches(1.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, 1.1, y + 0.1, 11, 0.5, title, font_size=22, color=color, bold=True)
    add_text(slide, 1.1, y + 0.6, 11, 1.0, desc, font_size=16, color=RGBColor(0xCC, 0xCC, 0xCC))

# ============================================================
# SLIDE 3: The Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'The Solution — GenBI Analytics Platform')

add_text(slide, 0.8, 1.2, 11.5, 0.8, 'Combines interactive BI dashboards with an AI-powered chatbot that answers\nbusiness questions in plain English — with full data lineage for every answer.', font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC))

capabilities = [
    ('Natural Language Q&A', 'Ask questions like "What is total\nrevenue by region?" and get\nSQL-backed answers instantly'),
    ('5 Interactive Dashboards', 'Executive, Sales & Menu,\nOperations, Customer Intelligence,\nFinancial Performance'),
    ('End-to-End Data Lineage', 'Every answer traces:\nSource → S3 → Glue ETL →\nRedshift → Dashboard'),
    ('Smart Dashboard Routing', 'Chatbot recommends the most\nrelevant dashboard for each\nquestion automatically'),
]

for i, (title, desc) in enumerate(capabilities):
    x = 0.8 + i * 3.1
    add_rounded_rect(slide, x, 2.3, 2.9, 3.2, RGBColor(0x1A, 0x2A, 0x3E))
    add_text(slide, x + 0.2, 2.5, 2.5, 0.6, title, font_size=17, color=ORANGE, bold=True)
    add_text(slide, x + 0.2, 3.2, 2.5, 2.0, desc, font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))

add_text(slide, 0.8, 5.8, 11, 0.5, 'Topic Guardrails: Chatbot only responds to restaurant operations questions — off-topic queries are politely redirected.', font_size=14, color=MEDIUM_GRAY)

# ============================================================
# SLIDE 4: Data Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'Demo Data — ABC Restaurant Group')

# KPI cards
kpis = [
    ('200', 'Stores', 'HK Island, Kowloon,\nNew Territories'),
    ('27M+', 'Rows', 'Star schema in\nAmazon Redshift'),
    ('HK$1.09B', 'Revenue', 'Full year 2023\ntransactions'),
    ('7.5M', 'Orders', 'Across 5 channels\nand 7 payment types'),
]

for i, (val, label, sub) in enumerate(kpis):
    x = 0.8 + i * 3.1
    add_rounded_rect(slide, x, 1.2, 2.9, 1.8, RGBColor(0x1A, 0x2A, 0x3E))
    add_text(slide, x + 0.2, 1.3, 2.5, 0.7, val, font_size=36, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.2, 1.9, 2.5, 0.4, label, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.2, 2.3, 2.5, 0.6, sub, font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Fact tables
add_text(slide, 0.8, 3.3, 5, 0.5, '8 Fact Tables', font_size=20, color=WHITE, bold=True)
facts = [
    'fact_sales — 17.5M rows — transaction line items',
    'fact_inventory — 2.19M rows — daily stock levels',
    'fact_labor — 665K rows — employee shifts & costs',
    'fact_service_performance — 1.75M rows — service times',
    'fact_customer_feedback — 28K rows — CSAT, NPS, sentiment',
    'fact_loyalty — 2.49M rows — loyalty program activity',
    'fact_equipment — 10K rows — maintenance & downtime',
    'fact_financial — 2.4K rows — monthly store P&L',
]

tb = add_text(slide, 0.8, 3.8, 6, 3.5, facts[0], font_size=13, color=RGBColor(0xBB, 0xBB, 0xBB))
for f in facts[1:]:
    add_para(tb.text_frame, f, font_size=13, color=RGBColor(0xBB, 0xBB, 0xBB), space_before=4)

# Dim tables
add_text(slide, 7.5, 3.3, 5, 0.5, '7 Dimension Tables', font_size=20, color=WHITE, bold=True)
dims = [
    'dim_date — 365 days (2023 + HK holidays)',
    'dim_store — 200 stores across 3 HK regions',
    'dim_menu_item — 30 items in 8 categories',
    'dim_channel — 5 order channels',
    'dim_payment_method — 7 payment types',
    'dim_promotion — 12 seasonal promotions',
    'dim_customer — 50,000 loyalty members',
]

tb2 = add_text(slide, 7.5, 3.8, 5.5, 3.0, dims[0], font_size=13, color=RGBColor(0xBB, 0xBB, 0xBB))
for d in dims[1:]:
    add_para(tb2.text_frame, d, font_size=13, color=RGBColor(0xBB, 0xBB, 0xBB), space_before=4)

# ============================================================
# SLIDE 5: Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'Solution Architecture')

# Pipeline flow - simplified visual
stages = [
    ('Source\nSystems', 'POS, HR, CRM,\nInventory, Finance', MEDIUM_GRAY),
    ('Amazon S3', 'Raw Landing Zone\nCSV partitions', GREEN_ACCENT),
    ('AWS Glue', '3 PySpark ETL\nJobs', BLUE),
    ('Amazon\nRedshift', 'genbi_mart\nStar Schema', RGBColor(0x8E, 0x44, 0xAD)),
    ('Amazon\nBedrock', 'KB (RAG) +\nLLM (Text-to-SQL)', ORANGE),
    ('Amazon\nQuickSight', '5 Embedded\nDashboards', RGBColor(0x29, 0x80, 0xB9)),
]

for i, (name, desc, color) in enumerate(stages):
    x = 0.5 + i * 2.15
    add_rounded_rect(slide, x, 1.3, 1.95, 1.1, color, name, font_size=13, font_color=WHITE, bold=True)
    add_text(slide, x, 2.5, 1.95, 0.6, desc, font_size=11, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)
    # Arrow
    if i < len(stages) - 1:
        add_text(slide, x + 1.85, 1.55, 0.4, 0.5, '→', font_size=24, color=ORANGE, bold=True)

# Presentation layer
add_text(slide, 0.8, 3.4, 11, 0.5, 'Presentation Layer', font_size=20, color=WHITE, bold=True)

pres_items = [
    ('Amazon CloudFront', 'CDN + HTTPS delivery'),
    ('Amazon API Gateway', 'REST API routing'),
    ('AWS Lambda', 'Serverless compute\n(chat + embed URLs)'),
    ('Amazon S3', 'Static website hosting'),
    ('End User Browser', 'Dashboard + Chatbot'),
]

for i, (name, desc) in enumerate(pres_items):
    x = 0.5 + i * 2.55
    add_rounded_rect(slide, x, 3.9, 2.35, 0.8, RGBColor(0x1A, 0x3A, 0x5F), name, font_size=12, font_color=LIGHT_BLUE, bold=True)
    add_text(slide, x, 4.75, 2.35, 0.5, desc, font_size=11, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Bedrock KB detail
add_text(slide, 0.8, 5.5, 12, 0.5, 'Knowledge Base: 6 indexed documents — schema overview, data lineage, SQL examples, business glossary, dashboard catalog, pipeline lineage', font_size=13, color=MEDIUM_GRAY)
add_text(slide, 0.8, 5.9, 12, 0.5, 'Model-agnostic: Amazon Bedrock supports Claude, Titan, Llama, Mistral, Cohere — swap model with one config change', font_size=13, color=MEDIUM_GRAY)

# ============================================================
# SLIDE 6: How the Chatbot Works
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'How the GenBI Chatbot Works')

steps = [
    ('1', 'User Asks Question', '"What is total revenue by region?"', BLUE),
    ('2', 'RAG Retrieval', 'Bedrock KB retrieves schema,\nlineage, SQL examples from\nOpenSearch Serverless', GREEN_ACCENT),
    ('3', 'SQL Generation', 'Bedrock LLM generates\nRedshift SQL with proper\nJOINs and aggregations', ORANGE),
    ('4', 'Query Execution', 'Redshift Data API executes\nSQL and returns results\nin ~2-3 seconds', RGBColor(0x8E, 0x44, 0xAD)),
    ('5', 'Response', 'Results table + explanation +\ndata lineage + recommended\ndashboard + SQL', LIGHT_BLUE),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.3 + i * 2.6
    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.7), Inches(1.3), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(18)
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, x + 0.1, 2.0, 2.4, 0.5, title, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.1, 2.5, 2.4, 1.5, desc, font_size=12, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

    if i < len(steps) - 1:
        add_text(slide, x + 2.35, 1.35, 0.4, 0.5, '→', font_size=24, color=ORANGE, bold=True)

# Data lineage example
add_text(slide, 0.8, 4.3, 11, 0.5, 'Data Lineage Example (included in every response):', font_size=16, color=WHITE, bold=True)
lineage_box = add_rounded_rect(slide, 0.8, 4.8, 11.7, 1.6, RGBColor(0x1A, 0x2A, 0x3E))
add_text(slide, 1.0, 4.9, 11.3, 1.5,
    'Source System: POS terminals (store registers)  →  S3 Raw: s3://.../pos/transactions/ + pos/line_items/  →\n'
    'Glue ETL (load_fact_sales): JOIN transactions + line_items ON transaction_id, calculate gross_profit  →\n'
    'Redshift genbi_mart.fact_sales: 17.5M rows, grain = one row per line item  →\n'
    'Dashboard Aggregation: COUNT(DISTINCT transaction_id) to convert line-item grain to order count',
    font_size=13, color=RGBColor(0xBB, 0xBB, 0xBB))

# ============================================================
# SLIDE 7: 5 Dashboards
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, '5 Interactive QuickSight Dashboards')

dashboards = [
    ('Executive Summary', 'Total revenue (HK$1.09B), orders (7.5M),\ngross profit, regional comparison,\nmonthly trends', 'fact_sales + dim_date + dim_store'),
    ('Sales & Menu', 'Revenue by category, top menu items,\nchannel mix, payment methods,\nhourly sales patterns', 'fact_sales + dim_menu_item +\ndim_channel + dim_payment_method'),
    ('Operations', 'Labor costs, staffing efficiency,\nshift productivity, workforce utilization,\nhours per shift', 'fact_labor + dim_store + dim_date'),
    ('Customer Intelligence', 'CSAT ratings, NPS scores,\nsentiment analysis, recommendation rate,\nrating by region', 'fact_customer_feedback +\ndim_store + dim_date'),
    ('Financial Performance', 'EBITDA, net profit margins,\ncost structure breakdown,\nmonthly P&L statements', 'fact_financial + dim_store +\ndim_date'),
]

for i, (name, metrics, sources) in enumerate(dashboards):
    x = 0.3 + i * 2.6
    add_rounded_rect(slide, x, 1.2, 2.4, 0.7, BLUE, name, font_size=13, font_color=WHITE, bold=True)
    add_text(slide, x + 0.1, 2.1, 2.3, 1.5, metrics, font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))
    add_text(slide, x + 0.1, 3.6, 2.3, 0.4, 'Data Sources:', font_size=10, color=ORANGE, bold=True)
    add_text(slide, x + 0.1, 3.9, 2.3, 1.0, sources, font_size=10, color=MEDIUM_GRAY)

add_text(slide, 0.8, 5.2, 11, 0.8,
    'All dashboards are embedded via Amazon QuickSight SDK with a unified chat sidebar.\n'
    'The chatbot recommends the most relevant dashboard for each question.',
    font_size=14, color=MEDIUM_GRAY)

# ============================================================
# SLIDE 8: Demo vs Enterprise Deployment
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'Deployment — Demo vs Enterprise')

# Demo column
add_rounded_rect(slide, 0.8, 1.2, 5.5, 0.6, BLUE, 'Demo Version (Current)', font_size=18, font_color=WHITE, bold=True)
demo_items = [
    'Public access via Amazon CloudFront',
    'Serverless: AWS Lambda + Amazon API Gateway',
    'No authentication (open access)',
    'Public endpoints for all services',
    'Near-zero cost when idle',
    'Ideal for: PoC, stakeholder demos',
]
tb = add_text(slide, 1.0, 2.0, 5.2, 3.5, demo_items[0], font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))
for item in demo_items[1:]:
    add_para(tb.text_frame, item, font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB), space_before=8)

# Enterprise column
add_rounded_rect(slide, 7.0, 1.2, 5.5, 0.6, ORANGE, 'Enterprise Version (Recommended)', font_size=18, font_color=WHITE, bold=True)
ent_items = [
    'Dual-path access: external (CloudFront + WAF) + internal (DX / VPN)',
    'Amazon Cognito / SAML SSO authentication',
    'VPC private subnets — no public endpoints',
    'VPC Endpoints for Bedrock & Redshift',
    'EC2 / ECS Fargate with Multi-AZ Auto Scaling',
    'AWS WAF, KMS, CloudTrail, X-Ray, QuickSight RLS',
]
tb2 = add_text(slide, 7.2, 2.0, 5.2, 3.5, ent_items[0], font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))
for item in ent_items[1:]:
    add_para(tb2.text_frame, item, font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB), space_before=8)

# Internal connectivity
add_text(slide, 0.8, 5.2, 12, 0.5, 'Internal Connectivity Options:', font_size=16, color=WHITE, bold=True)
conn = [
    ('AWS Direct Connect', 'HQ / regional offices — dedicated fiber, lowest latency'),
    ('AWS Site-to-Site VPN', 'Branch offices / DX backup — encrypted IPsec tunnel'),
    ('AWS Client VPN', 'Remote workers / WFH — per-user OpenVPN access'),
]
for i, (name, desc) in enumerate(conn):
    x = 0.8 + i * 4.2
    add_rounded_rect(slide, x, 5.7, 3.9, 0.7, RGBColor(0x1A, 0x3A, 0x5F))
    add_text(slide, x + 0.15, 5.72, 3.6, 0.35, name, font_size=13, color=LIGHT_BLUE, bold=True)
    add_text(slide, x + 0.15, 6.05, 3.6, 0.35, desc, font_size=11, color=MEDIUM_GRAY)

# ============================================================
# SLIDE 9: Tool-Agnostic Design
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'Tool-Agnostic & Model-Agnostic Design')

add_text(slide, 0.8, 1.2, 11, 0.8, 'The value is in the metadata documentation, not the specific tools.\nAny organization can enable the same AI Q&A experience over their data.', font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC))

replacements = [
    ('This Demo Uses', 'Can Be Replaced With', 'KB Docs Still Apply'),
    ('AWS Glue (ETL)', 'Informatica, dbt, Talend,\nAzure Data Factory, Airflow', 'data_lineage.md\npipeline_lineage.md'),
    ('Amazon Redshift (DW)', 'Snowflake, BigQuery,\nAzure Synapse, Databricks SQL', 'schema_overview.md\nsql_examples.md'),
    ('Amazon QuickSight (BI)', 'Power BI, Tableau,\nLooker, Apache Superset', 'dashboard_catalog.md'),
    ('Amazon Bedrock (LLM)', 'Azure OpenAI, Google Vertex AI,\nSelf-hosted models', 'Swap API call —\nKB + prompt unchanged'),
]

# Header row
for j, header in enumerate(replacements[0]):
    x = 0.8 + j * 4.1
    add_rounded_rect(slide, x, 2.3, 3.9, 0.5, BLUE, header, font_size=14, font_color=WHITE, bold=True)

# Data rows
for i, (col1, col2, col3) in enumerate(replacements[1:]):
    y = 2.95 + i * 1.05
    bg_color = RGBColor(0x1A, 0x2A, 0x3E) if i % 2 == 0 else RGBColor(0x1E, 0x2E, 0x42)
    for j, text in enumerate([col1, col2, col3]):
        x = 0.8 + j * 4.1
        add_rounded_rect(slide, x, y, 3.9, 0.9, bg_color, text, font_size=12, font_color=RGBColor(0xBB, 0xBB, 0xBB))

# ============================================================
# SLIDE 10: Detailed Data Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'Technical Deep Dive — End-to-End Data Pipeline')

# Row 1: Source Systems
add_text(slide, 0.5, 1.1, 2, 0.4, 'SOURCE SYSTEMS', font_size=12, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
src_systems = [('POS\nTerminals', 'Sales txns,\nline items'), ('HR /\nWorkforce', 'Labor shifts,\nschedules'), ('Inventory\nMgmt', 'Stock levels,\nwaste'), ('CRM /\nLoyalty', 'Profiles,\nfeedback'), ('Finance\n/ ERP', 'Monthly\nP&L')]
for i, (name, desc) in enumerate(src_systems):
    x = 0.3 + i * 2.55
    add_rounded_rect(slide, x, 1.5, 2.3, 0.7, MEDIUM_GRAY, name, font_size=10, font_color=WHITE, bold=True)
    add_text(slide, x, 2.2, 2.3, 0.4, desc, font_size=9, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)

# Arrow down
add_text(slide, 6.2, 2.5, 1, 0.4, 'CSV exports', font_size=11, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Row 2: S3 Raw Landing
add_text(slide, 0.5, 2.9, 2, 0.3, 'AMAZON S3', font_size=12, color=ORANGE, bold=True)
s3_paths = [
    ('pos/', 'pos_transactions_YYYY_MM.csv\npos_line_items_YYYY_MM.csv'),
    ('operations/', 'inventory_daily.csv\nlabor_shifts.csv\nservice_times.csv\nequipment_maintenance.csv'),
    ('customer/', 'customer_profiles.csv\ncustomer_feedback.csv\nloyalty_transactions.csv'),
    ('financial/', 'store_pl_monthly.csv\ncompetitor_pricing.csv'),
    ('reference/', 'stores.csv, menu_items.csv\nchannels.csv, payments.csv\npromotions.csv'),
]
for i, (folder, files) in enumerate(s3_paths):
    x = 0.3 + i * 2.55
    add_rounded_rect(slide, x, 3.15, 2.3, 1.3, RGBColor(0x1A, 0x3A, 0x2A))
    add_text(slide, x + 0.1, 3.18, 2.1, 0.3, 's3://bucket/' + folder, font_size=10, color=GREEN_ACCENT, bold=True)
    add_text(slide, x + 0.1, 3.48, 2.1, 0.9, files, font_size=9, color=RGBColor(0xAA, 0xAA, 0xAA))

# Arrow down
add_text(slide, 6.2, 4.4, 1, 0.3, 'PySpark', font_size=11, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Row 3: Glue ETL
add_text(slide, 0.5, 4.7, 2, 0.3, 'AWS GLUE ETL', font_size=12, color=ORANGE, bold=True)
glue_jobs = [
    ('load_dimensions', '7 dim tables\nDirect S3 → Redshift\nType casting, dedup', '~2 min'),
    ('load_facts', '8 fact tables\nJOINs across sources\nCalculated fields\n(gross_profit, labor_cost)', '~15 min'),
    ('load_metadata', 'ETL registry\nColumn-level lineage\nData dictionary', '~1 min'),
]
for i, (name, desc, duration) in enumerate(glue_jobs):
    x = 0.3 + i * 4.3
    add_rounded_rect(slide, x, 5.0, 4.0, 1.3, RGBColor(0x1A, 0x2A, 0x3E))
    add_text(slide, x + 0.15, 5.02, 3.7, 0.3, name, font_size=13, color=BLUE, bold=True)
    add_text(slide, x + 0.15, 5.32, 2.8, 0.9, desc, font_size=10, color=RGBColor(0xBB, 0xBB, 0xBB))
    add_text(slide, x + 3.0, 5.02, 0.9, 0.3, duration, font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)

# Arrow and Redshift target
add_text(slide, 0.5, 6.4, 12, 0.4, 'JDBC write  →  Amazon Redshift Serverless: genbi_mart schema  |  7 dimension tables + 8 fact tables  |  27M+ rows  |  Star schema optimized for analytics', font_size=12, color=RGBColor(0x8E, 0x44, 0xAD))

# ============================================================
# SLIDE 11: Star Schema Diagram
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'Technical Deep Dive — Star Schema (genbi_mart)')

# Center: Fact tables
fact_tables = [
    ('fact_sales', '17.5M rows', 'sale_key, date_key, store_id,\nitem_id, channel_id, transaction_id,\nquantity, unit_price, line_total,\ncogs_amount, gross_profit'),
    ('fact_labor', '665K rows', 'labor_key, date_key, store_id,\nemployee_id, shift_start, shift_end,\nhours_worked, hourly_rate, labor_cost'),
    ('fact_financial', '2.4K rows', 'financial_key, month_key, store_id,\nrevenue, cogs, labor_cost,\nrent, utilities, ebitda, net_profit'),
    ('fact_customer_feedback', '28K rows', 'feedback_key, date_key, store_id,\noverall_rating, food_rating,\nsentiment, would_recommend'),
]

for i, (name, rows, cols) in enumerate(fact_tables):
    x = 0.3 + i * 3.25
    add_rounded_rect(slide, x, 1.2, 3.05, 2.5, RGBColor(0x2A, 0x1A, 0x3E))
    add_text(slide, x + 0.15, 1.25, 2.8, 0.35, name, font_size=14, color=RGBColor(0x8E, 0x44, 0xAD), bold=True)
    add_text(slide, x + 0.15, 1.55, 2.8, 0.25, rows, font_size=11, color=ORANGE)
    add_text(slide, x + 0.15, 1.85, 2.8, 1.8, cols, font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA))

# More fact tables (row 2)
more_facts = [
    ('fact_inventory', '2.19M', 'date_key, store_id, item_id, opening_stock,\nreceived, sold, closing_stock, waste_qty'),
    ('fact_service_performance', '1.75M', 'date_key, store_id, channel_id, hour,\navg_service_time, orders_served'),
    ('fact_loyalty', '2.49M', 'date_key, customer_id, store_id,\npoints_earned, points_redeemed, order_value'),
    ('fact_equipment', '10K', 'store_id, equipment_type, event_date,\nevent_type, downtime_hours, repair_cost'),
]

for i, (name, rows, cols) in enumerate(more_facts):
    x = 0.3 + i * 3.25
    add_rounded_rect(slide, x, 3.8, 3.05, 1.5, RGBColor(0x2A, 0x1A, 0x3E))
    add_text(slide, x + 0.15, 3.82, 2.8, 0.3, name + '  (' + rows + ')', font_size=12, color=RGBColor(0x8E, 0x44, 0xAD), bold=True)
    add_text(slide, x + 0.15, 4.15, 2.8, 1.1, cols, font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA))

# Dimension tables
add_text(slide, 0.5, 5.5, 3, 0.3, 'DIMENSION TABLES', font_size=14, color=GREEN_ACCENT, bold=True)
dim_tables = [
    ('dim_date', '365', 'date_key, full_date, month_name,\nquarter, day_of_week, is_holiday'),
    ('dim_store', '200', 'store_id, store_name, region,\ndistrict, store_type, open_date'),
    ('dim_menu_item', '30', 'item_id, item_name, category,\nunit_price, cogs, food_cost_pct'),
    ('dim_channel', '5', 'channel_id,\nchannel_name'),
    ('dim_payment', '7', 'payment_method_id,\npayment_name'),
    ('dim_customer', '50K', 'customer_id, join_date,\ntier, lifetime_orders'),
    ('dim_promotion', '12', 'promo_id, promo_name,\nstart_date, discount_pct'),
]

for i, (name, rows, cols) in enumerate(dim_tables):
    x = 0.3 + i * 1.85
    add_rounded_rect(slide, x, 5.85, 1.7, 1.3, RGBColor(0x1A, 0x2A, 0x1A))
    add_text(slide, x + 0.08, 5.87, 1.55, 0.25, name, font_size=10, color=GREEN_ACCENT, bold=True)
    add_text(slide, x + 0.08, 6.1, 1.55, 0.15, rows + ' rows', font_size=8, color=ORANGE)
    add_text(slide, x + 0.08, 6.25, 1.55, 0.85, cols, font_size=8, color=RGBColor(0xAA, 0xAA, 0xAA))

# JOIN rules
add_text(slide, 0.5, 5.45, 12, 0.3, 'JOIN Rule: Fact tables use NATURAL keys (store_id, item_id) — Dimension tables use SURROGATE keys (store_key, item_key) — JOIN on natural keys only', font_size=10, color=YELLOW_ACCENT)

# ============================================================
# SLIDE 12: Knowledge Base Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'Technical Deep Dive — Knowledge Base Architecture')

# KB architecture flow
add_text(slide, 0.8, 1.1, 11, 0.6, 'The Knowledge Base is the brain that gives the LLM accurate context for SQL generation.\n6 markdown documents are chunked, embedded, and indexed in Amazon OpenSearch Serverless.', font_size=15, color=RGBColor(0xCC, 0xCC, 0xCC))

# KB docs - detailed
kb_docs = [
    ('01_schema_overview.md', 'Redshift Schema', 'Table definitions, column names & types,\nprimary/foreign keys, row counts,\njoin rules (natural vs surrogate keys),\ngrain of each fact table', BLUE),
    ('02_data_lineage.md', 'ETL Lineage', 'How each calculated metric is derived:\ngross_profit = line_total - cogs - discount\nETL schedule, data freshness SLAs,\nknown data characteristics', GREEN_ACCENT),
    ('03_sql_examples.md', 'SQL Patterns', 'Pre-validated SQL query patterns:\nrevenue by region, labor cost per order,\nwaste rate by category, monthly P&L,\ntop items by margin', ORANGE),
    ('04_business_glossary.md', 'Business Context', 'Metric definitions: AOV, CSAT, NPS, EBITDA\nHK market context: regions, currencies,\npayment methods (Octopus, etc.),\ntax rules, minimum wage', YELLOW_ACCENT),
    ('05_dashboard_catalog.md', 'Dashboard Map', 'Dashboard names, IDs, visual descriptions,\nwhich metrics each dashboard shows,\nfact/dim tables per dashboard,\nquestion-to-dashboard routing guide', RGBColor(0x29, 0x80, 0xB9)),
    ('06_pipeline_lineage.md', 'Pipeline Trace', 'For EACH table: source system →\nS3 raw file paths & columns →\nGlue job name, JOINs, transforms →\nRedshift table & grain →\ndashboard aggregation functions', RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (filename, title, desc, color) in enumerate(kb_docs):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.2
    y = 1.9 + row * 2.7
    add_rounded_rect(slide, x, y, 3.95, 2.4, RGBColor(0x1A, 0x2A, 0x3E))
    # Color bar on left
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(2.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, x + 0.15, y + 0.05, 3.7, 0.3, filename, font_size=11, color=color, bold=True)
    add_text(slide, x + 0.15, y + 0.35, 3.7, 0.3, title, font_size=14, color=WHITE, bold=True)
    add_text(slide, x + 0.15, y + 0.7, 3.7, 1.6, desc, font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# Indexing detail
add_text(slide, 0.5, 7.0, 12, 0.3, 'Indexing: Amazon Titan Embeddings V2  |  Chunking: 300 tokens, 20% overlap  |  Vector Store: Amazon OpenSearch Serverless  |  Top-K retrieval: 5 chunks per query', font_size=11, color=MEDIUM_GRAY)

# ============================================================
# SLIDE 12a: KB Doc 01 — Schema Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'KB Document: 01_schema_overview.md')
add_text(slide, 0.8, 0.95, 8, 0.35, 'What the LLM learns: table structures, column types, join rules, and grain of every table', font_size=13, color=MEDIUM_GRAY)

# How it contributes
add_rounded_rect(slide, 9.5, 0.3, 3.3, 0.95, BLUE)
add_text(slide, 9.65, 0.33, 3.0, 0.3, 'Contributes to:', font_size=11, color=WHITE, bold=True)
add_text(slide, 9.65, 0.6, 3.0, 0.5, 'SQL Query generation —\ncorrect table/column names & JOINs', font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# Fact table example
add_rounded_rect(slide, 0.5, 1.4, 6.1, 2.9, RGBColor(0x1A, 0x2A, 0x3E))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(0.06), Inches(2.9))
bar.fill.solid()
bar.fill.fore_color.rgb = RGBColor(0x8E, 0x44, 0xAD)
bar.line.fill.background()
add_text(slide, 0.75, 1.45, 5.7, 0.3, 'fact_sales — Transaction-level sales data', font_size=14, color=RGBColor(0x8E, 0x44, 0xAD), bold=True)
add_text(slide, 0.75, 1.8, 5.7, 2.4,
    'Columns: sale_key, transaction_id, date_key, store_id,\n'
    'item_id, customer_id, channel_id, payment_id, promo_id,\n'
    'order_hour, quantity, unit_price, line_total,\n'
    'discount_amount, cogs_amount, gross_profit\n\n'
    'Key formulas:\n'
    '  line_total = unit_price x quantity\n'
    '  gross_profit = line_total - cogs_amount - discount_amount\n\n'
    'Note: customer_id is NULL for ~60% of transactions\n'
    '(non-loyalty customers)',
    font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# Dimension table example
add_rounded_rect(slide, 6.8, 1.4, 6.0, 2.9, RGBColor(0x1A, 0x2A, 0x3E))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.4), Inches(0.06), Inches(2.9))
bar.fill.solid()
bar.fill.fore_color.rgb = GREEN_ACCENT
bar.line.fill.background()
add_text(slide, 7.05, 1.45, 5.5, 0.3, 'dim_date — Date dimension (365 rows)', font_size=14, color=GREEN_ACCENT, bold=True)
add_text(slide, 7.05, 1.8, 5.5, 1.2,
    'Columns: date_key, full_date, year, quarter, month,\n'
    'month_name, week_of_year, day_of_month, day_of_week,\n'
    'day_name, is_weekend, is_holiday\n\n'
    'is_holiday flags 17 HK public holidays:\n'
    'Lunar New Year, Easter, HKSAR Day, National Day, Xmas',
    font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

add_text(slide, 7.05, 3.1, 5.5, 0.3, 'dim_menu_item — Menu items (30 rows)', font_size=14, color=GREEN_ACCENT, bold=True)
add_text(slide, 7.05, 3.4, 5.5, 0.8,
    'Categories: Burgers, Chicken, Sides, Beverages,\n'
    'Breakfast, Desserts, Value Meals, Limited Time\n'
    'food_cost_pct = cogs / unit_price x 100',
    font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# JOIN rules - critical section
add_rounded_rect(slide, 0.5, 4.5, 12.3, 2.5, RGBColor(0x2A, 0x1A, 0x1A))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(0.06), Inches(2.5))
bar.fill.solid()
bar.fill.fore_color.rgb = RED_ACCENT
bar.line.fill.background()
add_text(slide, 0.75, 4.55, 11.8, 0.3, 'CRITICAL: JOIN Key Rules (most common source of SQL errors)', font_size=14, color=RED_ACCENT, bold=True)
add_text(slide, 0.75, 4.9, 5.5, 2.0,
    'Fact tables use NATURAL keys:\n'
    '  store_id, item_id, channel_id,\n'
    '  payment_method_id, customer_id\n\n'
    'Dimension tables use SURROGATE keys:\n'
    '  store_key, item_key, channel_key,\n'
    '  payment_key (internal use only)',
    font_size=12, color=RGBColor(0xBB, 0xBB, 0xBB))
add_text(slide, 6.5, 4.9, 6.0, 2.0,
    'Correct JOINs (always use natural keys):\n'
    '  fact_sales f JOIN dim_store s ON f.store_id = s.store_id\n'
    '  fact_sales f JOIN dim_menu_item m ON f.item_id = m.item_id\n'
    '  fact_sales f JOIN dim_date d ON f.date_key = d.date_key\n\n'
    'NEVER use: f.store_key, f.item_key (these don\'t exist on facts)\n'
    'fix_generated_sql() auto-corrects this if LLM gets it wrong',
    font_size=12, color=RGBColor(0xBB, 0xBB, 0xBB))

# ============================================================
# SLIDE 12b: KB Doc 02 — Data Lineage
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'KB Document: 02_data_lineage.md')
add_text(slide, 0.8, 0.95, 8, 0.35, 'What the LLM learns: how every calculated metric is derived, step by step', font_size=13, color=MEDIUM_GRAY)

add_rounded_rect(slide, 9.5, 0.3, 3.3, 0.95, ORANGE)
add_text(slide, 9.65, 0.33, 3.0, 0.3, 'Contributes to:', font_size=11, color=WHITE, bold=True)
add_text(slide, 9.65, 0.6, 3.0, 0.5, 'Data Lineage — metric\ncalculation formulas & sources', font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# Revenue & Profitability
formulas = [
    ('Revenue & Profitability', [
        ('gross_profit', 'line_total - cogs_amount - discount_amount', 'Source: JOIN of pos_transactions + pos_line_items on transaction_id'),
        ('line_total', 'unit_price x quantity', 'From pos_line_items.line_total'),
        ('discount_amount', 'Transaction-level discount evenly split across line items', 'COALESCE(discount, 0.0) — null-safe'),
    ], RGBColor(0x8E, 0x44, 0xAD)),
    ('Inventory & Waste', [
        ('waste_rate', 'units_wasted / (units_sold + units_wasted) x 100', '3% means: for every 100 units consumed, 3 wasted. Target: below 3%'),
        ('closing_stock', 'opening_stock + received - sold - wasted', 'From inventory_daily.csv with item_id prefix "M" stripped'),
    ], GREEN_ACCENT),
    ('Financial P&L', [
        ('ebitda', 'gross_profit - labor_cost - rent - utilities - marketing - other_opex', 'From store_pl_monthly.csv'),
        ('net_profit', 'ebitda - depreciation', 'net_margin target: 10-15% for healthy stores'),
        ('gross_margin_pct', 'gross_profit / revenue x 100', 'Benchmark: 60-70% for QSR industry'),
    ], BLUE),
    ('Labor Productivity', [
        ('labor_cost', 'hours_worked x hourly_rate', 'From labor_shifts.csv — one row per employee per shift'),
        ('labor_cost_per_order', 'total_labor / COUNT(DISTINCT txn_id)', 'Requires JOIN fact_labor + fact_sales on store_id + date_key'),
    ], YELLOW_ACCENT),
]

for i, (section, metrics, color) in enumerate(formulas):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 1.4 + row * 3.0
    add_rounded_rect(slide, x, y, 6.1, 2.75, RGBColor(0x1A, 0x2A, 0x3E))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(2.75))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, x + 0.2, y + 0.05, 5.7, 0.3, section, font_size=14, color=color, bold=True)
    content = ''
    for name, formula, source in metrics:
        content += f'{name} = {formula}\n    {source}\n'
    add_text(slide, x + 0.2, y + 0.4, 5.7, 2.2, content.strip(), font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# ============================================================
# SLIDE 12c: KB Doc 03 — SQL Examples
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'KB Document: 03_sql_examples.md')
add_text(slide, 0.8, 0.95, 8, 0.35, 'What the LLM learns: 14 pre-validated SQL patterns — the LLM mimics these proven templates', font_size=13, color=MEDIUM_GRAY)

add_rounded_rect(slide, 9.5, 0.3, 3.3, 0.95, BLUE)
add_text(slide, 9.65, 0.33, 3.0, 0.3, 'Contributes to:', font_size=11, color=WHITE, bold=True)
add_text(slide, 9.65, 0.6, 3.0, 0.5, 'SQL Query — correct syntax,\nJOIN patterns, aggregations', font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# SQL example categories
sql_cats = [
    ('Revenue Queries', 'Total revenue by month\nRevenue by region\nTop 10 stores by revenue\nRevenue by category\nRevenue by channel\nRevenue by payment method', BLUE),
    ('Time Analysis', 'Hourly sales pattern\nHoliday vs non-holiday\ncomparison\nWeekend vs weekday\nperformance', GREEN_ACCENT),
    ('Operations', 'Service wait times\nby channel\nInventory waste rate\nby category\nLabor cost per order\nby region', ORANGE),
    ('Customer & Financial', 'CSAT satisfaction\nby region\nLoyalty tier distribution\nMonthly P&L summary\nStores with negative profit', RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (title, queries, color) in enumerate(sql_cats):
    x = 0.5 + i * 3.2
    add_rounded_rect(slide, x, 1.4, 3.0, 2.2, RGBColor(0x1A, 0x2A, 0x3E))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.4), Inches(3.0), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, x + 0.15, 1.55, 2.7, 0.3, title, font_size=13, color=color, bold=True)
    add_text(slide, x + 0.15, 1.9, 2.7, 1.6, queries, font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# Actual SQL sample
add_text(slide, 0.5, 3.8, 6, 0.3, 'Sample: Top 10 Stores by Revenue', font_size=14, color=WHITE, bold=True)
add_rounded_rect(slide, 0.5, 4.1, 6.1, 3.0, RGBColor(0x0A, 0x14, 0x1E))
add_text(slide, 0.7, 4.15, 5.7, 2.9,
    'SELECT s.store_id, s.store_name,\n'
    '       s.district, s.region,\n'
    '       SUM(f.line_total) AS revenue,\n'
    '       SUM(f.gross_profit) AS profit,\n'
    '       COUNT(DISTINCT f.transaction_id)\n'
    '         AS orders\n'
    'FROM genbi_mart.fact_sales f\n'
    'JOIN genbi_mart.dim_store s\n'
    '  ON f.store_id = s.store_id\n'
    'GROUP BY s.store_id, s.store_name,\n'
    '         s.district, s.region\n'
    'ORDER BY revenue DESC LIMIT 10;',
    font_size=11, color=RGBColor(0x77, 0xDD, 0x77), font_name='Courier New')

add_text(slide, 6.8, 3.8, 6, 0.3, 'Sample: Inventory Waste Rate by Category', font_size=14, color=WHITE, bold=True)
add_rounded_rect(slide, 6.8, 4.1, 6.0, 3.0, RGBColor(0x0A, 0x14, 0x1E))
add_text(slide, 7.0, 4.15, 5.6, 2.9,
    'SELECT m.category_name,\n'
    '  ROUND(AVG(i.waste_rate), 2)\n'
    '    AS avg_waste_rate,\n'
    '  SUM(i.units_wasted) AS total_wasted,\n'
    '  SUM(i.units_sold) AS total_sold\n'
    'FROM genbi_mart.fact_inventory i\n'
    'JOIN genbi_mart.dim_menu_item m\n'
    '  ON i.item_id = m.item_id\n'
    'GROUP BY m.category_name\n'
    'ORDER BY avg_waste_rate DESC;',
    font_size=11, color=RGBColor(0x77, 0xDD, 0x77), font_name='Courier New')

# ============================================================
# SLIDE 12d: KB Doc 04 — Business Glossary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'KB Document: 04_business_glossary.md')
add_text(slide, 0.8, 0.95, 8, 0.35, 'What the LLM learns: business metric definitions & Hong Kong market context for accurate interpretation', font_size=13, color=MEDIUM_GRAY)

add_rounded_rect(slide, 9.5, 0.3, 3.3, 0.95, YELLOW_ACCENT)
add_text(slide, 9.65, 0.33, 3.0, 0.3, 'Contributes to:', font_size=11, color=WHITE, bold=True)
add_text(slide, 9.65, 0.6, 3.0, 0.5, 'Explanation & Assumptions —\nmetric meanings & benchmarks', font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# Key metrics
metrics_def = [
    ('AOV (Average Order Value)', 'SUM(line_total) / COUNT(DISTINCT transaction_id)\nTypical range: HKD 40-80'),
    ('Food Cost %', 'COGS / selling price x 100\nStored in dim_menu_item.food_cost_pct\nIndustry target: 28-35%'),
    ('CSAT', 'Customer Satisfaction Score\n1-5 scale from fact_customer_feedback\nTarget: above 4.0'),
    ('NPS (Net Promoter)', 'would_recommend TRUE/FALSE\n% promoters - % detractors\nTarget: above +30'),
    ('Waste Rate', 'units_wasted / (sold + wasted) x 100\nTarget: below 3%'),
    ('EBITDA', 'revenue - cogs - labor - rent -\nutilities - marketing - other_opex\nBenchmark: 15-20% of revenue'),
]

for i, (name, definition) in enumerate(metrics_def):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.2
    y = 1.4 + row * 1.75
    add_rounded_rect(slide, x, y, 3.95, 1.55, RGBColor(0x1A, 0x2A, 0x3E))
    add_text(slide, x + 0.15, y + 0.05, 3.65, 0.3, name, font_size=13, color=YELLOW_ACCENT, bold=True)
    add_text(slide, x + 0.15, y + 0.4, 3.65, 1.1, definition, font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# HK Market Context
add_text(slide, 0.5, 4.95, 12, 0.3, 'Hong Kong Market Context (embedded in LLM understanding)', font_size=14, color=WHITE, bold=True)

hk_items = [
    ('Regions', 'HK Island: CBD, high-rent, tourist/office\nKowloon: Dense residential, highest density\nNew Territories: Suburban, growing pop.', BLUE),
    ('Payment Methods', 'Octopus card: ~35% (dominant)\nCash: ~15% (declining)\nMobile pay: Apple/Alipay/WeChat (growing)\nVisa/Mastercard: ~25%', GREEN_ACCENT),
    ('Order Channels', 'Kiosk: ~35% (highest share)\nCounter: ~30% (traditional)\nMobile App: ~15% (order ahead)\nDelivery: ~12% | Drive-thru: ~8%', ORANGE),
]

for i, (title, desc, color) in enumerate(hk_items):
    x = 0.5 + i * 4.2
    add_rounded_rect(slide, x, 5.3, 3.95, 1.8, RGBColor(0x1A, 0x2A, 0x3E))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(5.3), Inches(0.06), Inches(1.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, x + 0.2, 5.35, 3.6, 0.3, title, font_size=13, color=color, bold=True)
    add_text(slide, x + 0.2, 5.65, 3.6, 1.4, desc, font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# ============================================================
# SLIDE 12e: KB Doc 05 — Dashboard Catalog
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'KB Document: 05_dashboard_catalog.md')
add_text(slide, 0.8, 0.95, 8, 0.35, 'What the LLM learns: which dashboards exist, what they show, and when to recommend each one', font_size=13, color=MEDIUM_GRAY)

add_rounded_rect(slide, 9.5, 0.3, 3.3, 0.95, RGBColor(0x29, 0x80, 0xB9))
add_text(slide, 9.65, 0.33, 3.0, 0.3, 'Contributes to:', font_size=11, color=WHITE, bold=True)
add_text(slide, 9.65, 0.6, 3.0, 0.5, 'Dashboard Recommendation —\nroutes user to right visual', font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# Dashboard specs
dash_specs = [
    ('Executive Summary', 'genbi-exec-dashboard', 'Senior leadership, regional managers',
     'Total Revenue KPI: SUM(fact_sales.line_total)\n'
     'Total Orders KPI: COUNT(DISTINCT transaction_id)\n'
     'Gross Profit KPI: SUM(gross_profit)\n'
     'Revenue by Region bar chart (HK Island/Kowloon/NT)\n'
     'Monthly Revenue Trend line chart (Jan-Dec 2023)\n'
     'Orders by Store Type pie chart'),
    ('Sales & Menu Analytics', 'genbi-sales-dashboard', 'Menu planners, marketing',
     'Revenue by Category bar chart (8 categories)\n'
     'Top Menu Items by Revenue bar chart\n'
     'Revenue by Channel pie chart (5 channels)\n'
     'Revenue by Payment Method pie chart (7 types)\n'
     'Sales by Hour of Day line chart (daily curve)'),
    ('Customer Intelligence', 'genbi-cust-dashboard', 'CX team, store managers',
     'Avg Customer Rating KPI (1-5 scale)\n'
     'Avg Rating by Region bar chart\n'
     'Sentiment by Quarter stacked bar\n'
     '  (positive / neutral / negative)\n'
     'Would Recommend pie chart'),
]

for i, (name, dash_id, audience, visuals) in enumerate(dash_specs):
    x = 0.5 + i * 4.2
    add_rounded_rect(slide, x, 1.4, 3.95, 3.7, RGBColor(0x1A, 0x2A, 0x3E))
    add_text(slide, x + 0.15, 1.45, 3.65, 0.35, name, font_size=14, color=RGBColor(0x29, 0x80, 0xB9), bold=True)
    add_text(slide, x + 0.15, 1.8, 3.65, 0.25, 'ID: ' + dash_id, font_size=10, color=ORANGE)
    add_text(slide, x + 0.15, 2.05, 3.65, 0.25, 'Audience: ' + audience, font_size=10, color=MEDIUM_GRAY)
    add_text(slide, x + 0.15, 2.35, 3.65, 2.6, visuals, font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# Recommendation routing table
add_text(slide, 0.5, 5.3, 12, 0.3, 'Dashboard Recommendation Guide (embedded in system prompt)', font_size=14, color=WHITE, bold=True)

routes = [
    ('Total revenue, overall perf.', 'Executive Summary'),
    ('Menu items, best sellers', 'Sales & Menu'),
    ('Channels (kiosk, delivery)', 'Sales & Menu'),
    ('Labor costs, staffing', 'Operations'),
    ('Customer ratings, CSAT', 'Customer Intelligence'),
    ('EBITDA, net profit, P&L', 'Financial Performance'),
]

for i, (topic, dashboard) in enumerate(routes):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.2
    y = 5.7 + row * 0.65
    add_rounded_rect(slide, x, y, 3.95, 0.55, RGBColor(0x1A, 0x2A, 0x3E))
    add_text(slide, x + 0.15, y + 0.02, 2.2, 0.5, topic, font_size=10, color=RGBColor(0xBB, 0xBB, 0xBB))
    add_text(slide, x + 2.4, y + 0.02, 1.4, 0.5, dashboard, font_size=10, color=ORANGE, bold=True, alignment=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 12f: KB Doc 06 — Pipeline Lineage
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'KB Document: 06_pipeline_lineage.md')
add_text(slide, 0.8, 0.95, 8, 0.35, 'What the LLM learns: complete source-to-dashboard trace for every table — how raw numbers transform', font_size=13, color=MEDIUM_GRAY)

add_rounded_rect(slide, 9.5, 0.3, 3.3, 0.95, RGBColor(0x8E, 0x44, 0xAD))
add_text(slide, 9.65, 0.33, 3.0, 0.3, 'Contributes to:', font_size=11, color=WHITE, bold=True)
add_text(slide, 9.65, 0.6, 3.0, 0.5, 'Data Lineage — end-to-end\npipeline trace in every answer', font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# fact_sales full pipeline
add_text(slide, 0.5, 1.4, 12, 0.3, 'Example: fact_sales — How "Total Orders" Goes from Register to Dashboard', font_size=15, color=WHITE, bold=True)

pipeline_stages = [
    ('1. Source System', 'POS Terminals', 'Raw transaction data\ngenerated at point of sale.\nOne record per transaction +\none record per line item.', MEDIUM_GRAY),
    ('2. S3 Raw Landing', 's3://.../pos/', 'pos_transactions_YYYY_MM.csv\n  (transaction_id, store_id,\n   order_date, channel_id)\n\npos_line_items_YYYY_MM.csv\n  (transaction_id, item_id,\n   quantity, unit_price, line_total)', GREEN_ACCENT),
    ('3. Glue ETL', 'load_fact_sales', 'JOIN txns + items\n  ON transaction_id (INNER)\n\nTransforms:\n  date_key = FORMAT(order_date)\n  discount = COALESCE(0.0)\n  gross_profit = line_total\n    - discount - cogs', BLUE),
    ('4. Redshift', 'genbi_mart.fact_sales', '17.5M rows loaded\nGrain: one row per line item\nper transaction\n\nStored as: sale_key (PK),\ntransaction_id, date_key,\nstore_id, item_id, ...', RGBColor(0x8E, 0x44, 0xAD)),
    ('5. Dashboard', 'Aggregation', 'COUNT(DISTINCT\n  transaction_id)\n\nConverts line-item grain\nto unique order count.\n\n17.5M line items →\n7.5M distinct orders', ORANGE),
]

for i, (stage, name, desc, color) in enumerate(pipeline_stages):
    x = 0.3 + i * 2.6
    add_rounded_rect(slide, x, 1.8, 2.35, 0.5, color, stage, font_size=11, font_color=WHITE, bold=True)
    add_rounded_rect(slide, x, 2.35, 2.35, 0.4, RGBColor(0x1A, 0x2A, 0x3E), name, font_size=10, font_color=color, bold=True)
    add_text(slide, x + 0.1, 2.8, 2.2, 2.5, desc, font_size=10, color=RGBColor(0xBB, 0xBB, 0xBB))
    if i < 4:
        add_text(slide, x + 2.25, 1.95, 0.4, 0.4, '→', font_size=18, color=ORANGE, bold=True)

# Second example: dim_menu_item pipeline
add_text(slide, 0.5, 5.1, 12, 0.3, 'Example: dim_menu_item — How Food Cost % Is Calculated', font_size=15, color=WHITE, bold=True)
add_rounded_rect(slide, 0.5, 5.4, 12.3, 1.7, RGBColor(0x1A, 0x2A, 0x3E))
add_text(slide, 0.7, 5.45, 11.8, 1.6,
    'Source: s3://.../reference/menu_items/ + reference/categories/\n'
    '  → menu_items.csv has: item_id, item_name, category_id, unit_price, cogs\n'
    '  → categories.csv has: category_id, category_name\n\n'
    'Glue ETL (load_dim_menu_item): JOIN menu_items + categories ON category_id\n'
    '  → Calculate: food_cost_pct = cogs / unit_price x 100\n'
    '  → Example: Big Mac: cogs = HK$12, price = HK$42 → food_cost_pct = 28.6%\n\n'
    'Target Redshift: genbi_mart.dim_menu_item (30 rows) → Used by QuickSight "Sales & Menu" dashboard for profitability analysis',
    font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# ============================================================
# SLIDE 13: RAG Retrieval Mechanics
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'Technical Deep Dive — RAG Retrieval Mechanics')

# 3-step RAG flow
rag_steps = [
    ('Step 1\nVector Search', 'User question → vector embedding →\ncompared against pre-indexed KB chunks\nin Amazon OpenSearch Serverless.\n\nTop-5 most relevant chunks retrieved.\n\nExample: "total orders" retrieves:\n  • fact_sales schema (join rules)\n  • COUNT(DISTINCT transaction_id) pattern\n  • POS pipeline lineage', BLUE),
    ('Step 2\nPrompt Assembly', 'Retrieved KB chunks injected into\nLLM prompt alongside:\n\n  • User\'s natural language question\n  • System instructions (SQL rules,\n    JOIN key rules, lineage format)\n  • Topic guardrails\n\nGives LLM precise, verified context\ninstead of general training knowledge.', GREEN_ACCENT),
    ('Step 3\nStructured Generation', 'LLM generates JSON response with\n5 fields, each from different KB docs:\n\n  • SQL query  ← schema + sql_examples\n  • Explanation ← business_glossary\n  • Dashboard  ← dashboard_catalog\n  • Lineage    ← pipeline_lineage\n  • Assumptions ← business_glossary\n\nSQL post-processed by fix_generated_sql()\nto correct common LLM mistakes.', ORANGE),
]

for i, (title, desc, color) in enumerate(rag_steps):
    x = 0.5 + i * 4.2
    add_rounded_rect(slide, x, 1.2, 3.95, 4.5, RGBColor(0x1A, 0x2A, 0x3E))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.2), Inches(3.95), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, x + 0.2, 1.4, 3.5, 0.6, title, font_size=16, color=color, bold=True)
    add_text(slide, x + 0.2, 2.1, 3.5, 3.5, desc, font_size=12, color=RGBColor(0xBB, 0xBB, 0xBB))

    if i < 2:
        add_text(slide, x + 3.85, 3.0, 0.5, 0.5, '→', font_size=28, color=ORANGE, bold=True)

# Response mapping table
add_text(slide, 0.5, 5.9, 12, 0.3, 'Response Field → KB Document Mapping', font_size=14, color=WHITE, bold=True)
mappings = [
    ('SQL Query', '01_schema_overview + 03_sql_examples'),
    ('Explanation', '04_business_glossary'),
    ('Dashboard', '05_dashboard_catalog'),
    ('Data Lineage', '06_pipeline_lineage + 02_data_lineage'),
    ('Assumptions', '04_business_glossary'),
]
for i, (field, docs) in enumerate(mappings):
    x = 0.5 + i * 2.6
    add_rounded_rect(slide, x, 6.25, 2.4, 0.45, RGBColor(0x1A, 0x3A, 0x5F), field, font_size=11, font_color=LIGHT_BLUE, bold=True)
    add_text(slide, x, 6.72, 2.4, 0.35, docs, font_size=9, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 14: Data Lineage Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'Technical Deep Dive — End-to-End Data Lineage')

add_text(slide, 0.8, 1.1, 11, 0.5, 'For every chatbot answer, the full pipeline is traced — from source system to the final aggregated metric.', font_size=15, color=RGBColor(0xCC, 0xCC, 0xCC))

# Lineage pipeline stages (horizontal flow)
lineage_stages = [
    ('Source System', 'POS terminals\n(store registers)', 'Raw transaction\ndata generated\nat point of sale', MEDIUM_GRAY),
    ('S3 Raw Landing', 's3://.../pos/\ntransactions/\n+ line_items/', 'Monthly CSV\npartitions\nYYYY_MM format', GREEN_ACCENT),
    ('Glue ETL', 'load_fact_sales\nJOIN txns + items\nON transaction_id', 'Calculate:\ngross_profit =\nline_total - cogs\n- discount', BLUE),
    ('Redshift', 'genbi_mart.\nfact_sales\n17.5M rows', 'Grain: one row\nper line item\nper transaction', RGBColor(0x8E, 0x44, 0xAD)),
    ('Dashboard\nAggregation', 'COUNT(DISTINCT\ntransaction_id)', 'Converts line-item\ngrain to unique\norder count', ORANGE),
]

for i, (stage, detail, note, color) in enumerate(lineage_stages):
    x = 0.3 + i * 2.6
    # Stage box
    add_rounded_rect(slide, x, 1.7, 2.35, 0.55, color, stage, font_size=12, font_color=WHITE, bold=True)
    add_rounded_rect(slide, x, 2.35, 2.35, 1.2, RGBColor(0x1A, 0x2A, 0x3E), detail, font_size=11, font_color=RGBColor(0xBB, 0xBB, 0xBB))
    add_text(slide, x, 3.6, 2.35, 0.8, note, font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
    if i < len(lineage_stages) - 1:
        add_text(slide, x + 2.25, 2.1, 0.4, 0.5, '→', font_size=22, color=ORANGE, bold=True)

# More lineage examples
add_text(slide, 0.8, 4.5, 11, 0.4, 'More Lineage Examples (from 06_pipeline_lineage.md):', font_size=16, color=WHITE, bold=True)

examples = [
    ('Total Revenue by Region', 'POS terminals → s3://.../pos/transactions/ + line_items/ →\nGlue (load_fact_sales): JOIN + calculate line_total →\nRedshift fact_sales JOIN dim_store → SUM(line_total) GROUP BY region'),
    ('Average CSAT Rating', 'Survey platform → s3://.../customer/customer_feedback.csv →\nGlue (load_facts): type cast ratings to DECIMAL →\nRedshift fact_customer_feedback → AVG(overall_rating)'),
    ('Monthly EBITDA', 'Finance/ERP → s3://.../financial/store_pl_monthly.csv →\nGlue (load_facts): direct load, month_key derivation →\nRedshift fact_financial → SUM(ebitda) GROUP BY month'),
    ('Labor Cost per Order', 'HR system → s3://.../operations/labor_shifts.csv + POS →\nGlue: JOIN labor to sales on store_id + date_key →\nRedshift fact_labor + fact_sales → labor_cost / COUNT(DISTINCT txn_id)'),
]

for i, (title, trace) in enumerate(examples):
    y = 4.95 + i * 0.65
    bg = RGBColor(0x1A, 0x2A, 0x3E) if i % 2 == 0 else RGBColor(0x1E, 0x2E, 0x42)
    add_rounded_rect(slide, 0.5, y, 12.3, 0.6, bg)
    add_text(slide, 0.7, y + 0.02, 2.5, 0.55, title, font_size=11, color=ORANGE, bold=True)
    add_text(slide, 3.3, y + 0.02, 9.3, 0.55, trace, font_size=10, color=RGBColor(0xBB, 0xBB, 0xBB))

# ============================================================
# SLIDE 15: SQL Generation & Post-Processing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'Technical Deep Dive — SQL Generation & Guardrails')

# Left: SQL generation flow
add_text(slide, 0.8, 1.1, 5.5, 0.4, 'SQL Generation Pipeline', font_size=18, color=WHITE, bold=True)

sql_flow = [
    ('System Prompt', 'Injected into every LLM call:\n• SQL rules (fully qualified table names)\n• JOIN key rules (natural vs surrogate)\n• Date filtering rules (dim_date.full_date)\n• COUNT(DISTINCT transaction_id) for orders\n• HKD currency, no sales tax\n• Lineage format template\n• Dashboard recommendation guide\n• Topic guardrails', BLUE),
    ('Post-Processing\n(fix_generated_sql)', 'Automatic corrections before execution:\n• Fix surrogate keys on fact tables:\n  store_key → store_id\n  item_key → item_id\n  channel_key → channel_id\n• Fix fact_financial date key:\n  date_key → month_key\n• Regex-based, alias-aware matching', ORANGE),
]

for i, (title, desc, color) in enumerate(sql_flow):
    y = 1.6 + i * 3.0
    add_rounded_rect(slide, 0.5, y, 5.8, 2.7, RGBColor(0x1A, 0x2A, 0x3E))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(0.06), Inches(2.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, 0.75, y + 0.1, 5.3, 0.4, title, font_size=14, color=color, bold=True)
    add_text(slide, 0.75, y + 0.55, 5.3, 2.1, desc, font_size=12, color=RGBColor(0xBB, 0xBB, 0xBB))

# Right: Topic guardrails & response format
add_text(slide, 7.0, 1.1, 5.5, 0.4, 'Guardrails & Response Format', font_size=18, color=WHITE, bold=True)

add_rounded_rect(slide, 6.8, 1.6, 6.0, 2.3, RGBColor(0x1A, 0x2A, 0x3E))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.6), Inches(0.06), Inches(2.3))
bar.fill.solid()
bar.fill.fore_color.rgb = RED_ACCENT
bar.line.fill.background()
add_text(slide, 7.05, 1.65, 5.5, 0.35, 'Topic Guardrails', font_size=14, color=RED_ACCENT, bold=True)
add_text(slide, 7.05, 2.0, 5.5, 1.8,
    'Allowed topics:\n'
    '  Sales, revenue, orders, menu items, inventory,\n'
    '  labor/staffing, service performance, customer\n'
    '  feedback, loyalty, equipment, financial P&L\n\n'
    'Off-topic → returns explanation with no SQL\n'
    '(politics, weather, coding, general knowledge)',
    font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

add_rounded_rect(slide, 6.8, 4.1, 6.0, 3.0, RGBColor(0x1A, 0x2A, 0x3E))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(4.1), Inches(0.06), Inches(3.0))
bar.fill.solid()
bar.fill.fore_color.rgb = GREEN_ACCENT
bar.line.fill.background()
add_text(slide, 7.05, 4.15, 5.5, 0.35, 'Structured JSON Response', font_size=14, color=GREEN_ACCENT, bold=True)
add_text(slide, 7.05, 4.5, 5.5, 2.5,
    '{\n'
    '  "sql": "SELECT ... FROM genbi_mart...",\n'
    '  "explanation": "This query calculates...",\n'
    '  "lineage": "Source System: POS → S3 → ...",\n'
    '  "assumptions": "Revenue = line_total...",\n'
    '  "recommended_dashboard": "genbi-exec-dashboard"\n'
    '}\n\n'
    'Parsed with json.JSONDecoder.raw_decode()\n'
    'for robust extraction from LLM output.',
    font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))

# ============================================================
# SLIDE 16: Enterprise Architecture Diagram
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'Enterprise Architecture — Dual-Path Access')

# External path (left)
add_text(slide, 0.3, 1.1, 6, 0.35, 'EXTERNAL ACCESS PATH', font_size=14, color=LIGHT_BLUE, bold=True)
ext_items = [
    ('External User', RGBColor(0x44, 0x44, 0x44)),
    ('AWS WAF', RED_ACCENT),
    ('Amazon CloudFront', BLUE),
    ('Amazon Cognito\n(SAML / OIDC)', YELLOW_ACCENT),
    ('Public ALB', BLUE),
]
for i, (name, color) in enumerate(ext_items):
    y = 1.5 + i * 0.75
    add_rounded_rect(slide, 0.3, y, 2.8, 0.6, color, name, font_size=11, font_color=WHITE, bold=True)
    if i < len(ext_items) - 1:
        add_text(slide, 1.5, y + 0.55, 0.5, 0.25, '|', font_size=14, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Internal path (right side)
add_text(slide, 3.8, 1.1, 6, 0.35, 'INTERNAL ACCESS PATH', font_size=14, color=GREEN_ACCENT, bold=True)
int_items = [
    ('Internal User\n(Staff / Analysts)', RGBColor(0x44, 0x44, 0x44)),
    ('Corporate Network', MEDIUM_GRAY),
    ('AWS Direct Connect\n/ Site-to-Site VPN\n/ Client VPN', GREEN_ACCENT),
    ('AWS Transit Gateway', GREEN_ACCENT),
    ('Private ALB', GREEN_ACCENT),
]
for i, (name, color) in enumerate(int_items):
    y = 1.5 + i * 0.75
    add_rounded_rect(slide, 3.8, y, 2.8, 0.6, color, name, font_size=11 if i != 2 else 9, font_color=WHITE, bold=True)
    if i < len(int_items) - 1:
        add_text(slide, 5.0, y + 0.55, 0.5, 0.25, '|', font_size=14, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Center: VPC
add_text(slide, 7.2, 1.1, 6, 0.35, 'CUSTOMER VPC (PRIVATE SUBNETS)', font_size=14, color=ORANGE, bold=True)

# Arrows from both ALBs to EC2
add_text(slide, 2.95, 4.5, 0.5, 0.3, '→', font_size=20, color=ORANGE, bold=True)
add_text(slide, 6.5, 4.5, 0.5, 0.3, '→', font_size=20, color=ORANGE, bold=True)

vpc_items = [
    ('Amazon EC2 / AWS Fargate\n(GenBI API — Multi-AZ)', BLUE, 7.2, 1.5, 5.5, 0.8),
    ('VPC Endpoint\n(Bedrock)', RGBColor(0x1A, 0x3A, 0x5F), 7.2, 2.5, 2.6, 0.65),
    ('VPC Endpoint\n(Redshift)', RGBColor(0x1A, 0x3A, 0x5F), 10.1, 2.5, 2.6, 0.65),
    ('Amazon Bedrock\nKnowledge Bases', ORANGE, 7.2, 3.4, 2.6, 0.65),
    ('Amazon Bedrock\nLLM', ORANGE, 10.1, 3.4, 2.6, 0.65),
    ('Amazon Redshift\nServerless', RGBColor(0x8E, 0x44, 0xAD), 7.2, 4.3, 2.6, 0.65),
    ('Amazon QuickSight\n(Row-Level Security)', RGBColor(0x29, 0x80, 0xB9), 10.1, 4.3, 2.6, 0.65),
]

for name, color, x, y, w, h in vpc_items:
    add_rounded_rect(slide, x, y, w, h, color, name, font_size=11, font_color=WHITE, bold=True)

# Operational layer
add_text(slide, 7.2, 5.2, 5.5, 0.3, 'OPERATIONAL LAYER', font_size=12, color=MEDIUM_GRAY, bold=True)
ops = [('Amazon\nCloudWatch', 7.2), ('AWS\nX-Ray', 9.1), ('AWS\nCloudTrail', 11.0)]
for name, x in ops:
    add_rounded_rect(slide, x, 5.5, 1.7, 0.6, RGBColor(0x1A, 0x3A, 0x5F), name, font_size=10, font_color=MEDIUM_GRAY)

# Security controls footer
add_text(slide, 0.3, 6.3, 12.5, 0.5,
    'Security: VPC Endpoints (no internet)  |  AWS KMS encryption at rest  |  TLS in transit  |  QuickSight RLS  |  AWS WAF (OWASP rules)  |  CloudTrail audit logs  |  X-Ray tracing',
    font_size=11, color=MEDIUM_GRAY)

# ============================================================
# SLIDE 17: Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, 'Try It Yourself & Next Steps')

# Links
add_rounded_rect(slide, 0.8, 1.3, 5.5, 1.0, RGBColor(0x1A, 0x3A, 0x5F))
add_text(slide, 1.0, 1.35, 5, 0.4, 'Live Demo', font_size=16, color=ORANGE, bold=True)
add_text(slide, 1.0, 1.75, 5, 0.4, 'https://d1k3nghlesd8gk.cloudfront.net', font_size=14, color=LIGHT_BLUE)

add_rounded_rect(slide, 7.0, 1.3, 5.5, 1.0, RGBColor(0x1A, 0x3A, 0x5F))
add_text(slide, 7.2, 1.35, 5, 0.4, 'Source Code & Documentation', font_size=16, color=ORANGE, bold=True)
add_text(slide, 7.2, 1.75, 5, 0.4, 'https://github.com/danielpeggy/bi-report-chatbot', font_size=14, color=LIGHT_BLUE)

# Sample questions
add_text(slide, 0.8, 2.7, 11, 0.5, 'Sample Questions to Try:', font_size=20, color=WHITE, bold=True)
questions = [
    '"What is the total revenue by region?"',
    '"Which menu items have the highest profit margin?"',
    '"Show me monthly order trends for 2023"',
    '"What is the average CSAT rating by store?"',
    '"What are the labor costs per order by region?"',
]
tb = add_text(slide, 1.0, 3.2, 11, 2.0, questions[0], font_size=15, color=RGBColor(0xBB, 0xBB, 0xBB))
for q in questions[1:]:
    add_para(tb.text_frame, q, font_size=15, color=RGBColor(0xBB, 0xBB, 0xBB), space_before=8)

# Next steps
add_text(slide, 0.8, 5.0, 11, 0.5, 'Recommended Next Steps:', font_size=20, color=WHITE, bold=True)

steps = [
    ('1. Live Demo Walkthrough', '30-minute session to explore dashboards and chatbot capabilities'),
    ('2. Architecture Workshop', 'Map your current data pipeline to the GenBI approach'),
    ('3. PoC Scoping', 'Identify 2-3 dashboards and a target dataset for proof-of-concept'),
]

for i, (title, desc) in enumerate(steps):
    x = 0.8 + i * 4.2
    add_rounded_rect(slide, x, 5.5, 3.9, 1.2, RGBColor(0x1A, 0x3A, 0x5F))
    add_text(slide, x + 0.15, 5.55, 3.6, 0.4, title, font_size=14, color=ORANGE, bold=True)
    add_text(slide, x + 0.15, 5.95, 3.6, 0.5, desc, font_size=12, color=RGBColor(0xBB, 0xBB, 0xBB))


# Save
output = 'GenBI_Solution_Presentation.pptx'
prs.save(output)
print(f'Saved: {output}')
print(f'Slides: {len(prs.slides)}')
